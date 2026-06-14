"""Agent 5 — PPT generation. Renders a DeckPlan into a .pptx (python-pptx).

Returns the deck as bytes so the caller can hand it to any storage backend.
Uses templates/brand.pptx as the master if present.
"""
from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from app.agents.analysis_agent import DeckPlan, SlidePlan
from app.core.base import Agent

_BACKEND_ROOT = Path(__file__).resolve().parents[2]
_DEFAULT_TEMPLATE = _BACKEND_ROOT / "templates" / "brand.pptx"


def safe_filename(name: str) -> str:
    return re.sub(r"[^\w\- ]+", "", name).strip()[:80] or "deck"


class PptAgent(Agent):
    name = "ppt"

    def __init__(self, template: Path | None = None) -> None:
        tpl = template or _DEFAULT_TEMPLATE
        self.template = tpl if tpl and tpl.exists() else None

    # --- tool: build a pptx into memory ---
    def build(self, plan: DeckPlan) -> bytes:
        prs = Presentation(str(self.template)) if self.template else Presentation()
        for i, slide in enumerate(plan.slides):
            if i == 0:
                self._title_slide(prs, plan.deck_title, slide)
            else:
                self._content_slide(prs, slide)
        buf = BytesIO()
        prs.save(buf)
        data = buf.getvalue()
        self.log(f"rendered '{plan.deck_title}' — {len(plan.slides)} slides, "
                 f"{len(data)} bytes")
        return data

    # --- layouts ---
    def _title_slide(self, prs: Presentation, deck_title: str, slide: SlidePlan) -> None:
        s = prs.slides.add_slide(prs.slide_layouts[0])
        if s.shapes.title is not None:
            s.shapes.title.text = deck_title
        if len(s.placeholders) > 1:
            s.placeholders[1].text = (" | ".join(slide.bullets)
                                      if slide.bullets else slide.title)

    def _content_slide(self, prs: Presentation, slide: SlidePlan) -> None:
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        s = prs.slides.add_slide(layout)
        if s.shapes.title is not None:
            s.shapes.title.text = slide.title
        body = self._body_placeholder(s)
        if body is not None:
            tf = body.text_frame
            tf.word_wrap = True
            tf.clear()
            for j, bullet in enumerate(slide.bullets or [slide.title]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(18)
        if slide.notes:
            s.notes_slide.notes_text_frame.text = slide.notes

    @staticmethod
    def _body_placeholder(slide):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:  # not the title
                return ph
        return slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.5), Inches(5))